"""
本地文件解析器模块
当 MinerU API 不可用时的备用解析方案
"""
import logging
from pathlib import Path
from typing import Dict, Any

logger = logging.getLogger(__name__)


def parse_pdf_basic(file_path: str) -> Dict[str, Any]:
    """基础 PDF 解析（使用 pdfplumber）"""
    try:
        import pdfplumber
        
        text_blocks = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_blocks.append({
                        "content": page_text,
                        "page": i + 1
                    })
        
        full_text = "\n\n".join([b["content"] for b in text_blocks])
        
        return {
            "text_blocks": text_blocks,
            "images": [],
            "tables": [],
            "full_text": full_text
        }
        
    except Exception as e:
        logger.error(f"基础 PDF 解析失败 {file_path}: {e}")
        return {"text_blocks": [], "images": [], "tables": [], "full_text": ""}


def parse_docx(file_path: str) -> Dict[str, Any]:
    """解析 Word 文档"""
    from docx import Document
    
    text_blocks = []
    images = []
    tables = []
    
    try:
        doc = Document(file_path)
        
        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                text_blocks.append({"content": para.text, "page": 0})
        
        # 提取表格
        for table in doc.tables:
            table_markdown = _docx_table_to_markdown(table)
            tables.append({"markdown": table_markdown, "page": 0})
        
        full_text = _merge_content(text_blocks, images, tables)
        
        return {
            "text_blocks": text_blocks,
            "images": images,
            "tables": tables,
            "full_text": full_text
        }
        
    except Exception as e:
        logger.error(f"Word 文档解析失败 {file_path}: {e}")
        return {"text_blocks": [], "images": [], "tables": [], "full_text": ""}


def parse_pptx(file_path: str) -> Dict[str, Any]:
    """解析 PowerPoint 文档"""
    from pptx import Presentation
    
    text_blocks = []
    images = []
    tables = []
    
    try:
        prs = Presentation(file_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_blocks.append({
                        "content": shape.text,
                        "page": slide_idx + 1
                    })
                # 提取图片
                elif shape.shape_type == 13:  # PICTURE
                    images.append({
                        "image_data": shape.image.blob,
                        "ocr_text": "",
                        "page": slide_idx + 1
                    })
                # 提取表格
                elif shape.shape_type == 19:  # TABLE
                    table_markdown = _pptx_table_to_markdown(shape.table)
                    tables.append({
                        "markdown": table_markdown,
                        "page": slide_idx + 1
                    })
        
        full_text = _merge_content(text_blocks, images, tables)
        
        return {
            "text_blocks": text_blocks,
            "images": images,
            "tables": tables,
            "full_text": full_text
        }
        
    except Exception as e:
        logger.error(f"PowerPoint 文档解析失败 {file_path}: {e}")
        return {"text_blocks": [], "images": [], "tables": [], "full_text": ""}


def parse_xlsx(file_path: str) -> Dict[str, Any]:
    """解析 Excel 文档"""
    import pandas as pd
    
    text_blocks = []
    tables = []
    
    try:
        # 读取所有 sheet
        xls = pd.ExcelFile(file_path)
        
        for sheet_idx, sheet_name in enumerate(xls.sheet_names):
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # 转为 Markdown 表格
            table_markdown = df.to_markdown(index=False)
            tables.append({
                "markdown": table_markdown,
                "page": sheet_idx + 1
            })
            
            # 添加简要描述
            text_blocks.append({
                "content": f"工作表：{sheet_name}，共 {len(df)} 行数据",
                "page": sheet_idx + 1
            })
        
        full_text = _merge_content(text_blocks, [], tables)
        
        return {
            "text_blocks": text_blocks,
            "images": [],
            "tables": tables,
            "full_text": full_text
        }
        
    except Exception as e:
        logger.error(f"Excel 文档解析失败 {file_path}: {e}")
        return {"text_blocks": [], "images": [], "tables": [], "full_text": ""}


def parse_txt(file_path: str) -> Dict[str, Any]:
    """解析纯文本文件"""
    import chardet
    
    try:
        # 检测编码
        with open(file_path, "rb") as f:
            raw = f.read(10000)
        result = chardet.detect(raw)
        encoding = result.get("encoding", "utf-8") or "utf-8"
        
        # 读取内容
        with open(file_path, "r", encoding=encoding, errors="replace") as f:
            content = f.read()
        
        # 按段落分割
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        text_blocks = [{"content": p, "page": 0} for p in paragraphs]
        
        return {
            "text_blocks": text_blocks,
            "images": [],
            "tables": [],
            "full_text": content
        }
        
    except Exception as e:
        logger.error(f"文本文件解析失败 {file_path}: {e}")
        return {"text_blocks": [], "images": [], "tables": [], "full_text": ""}


def _docx_table_to_markdown(table) -> str:
    """将 Word 表格转为 Markdown"""
    lines = []
    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(lines)


def _pptx_table_to_markdown(table) -> str:
    """将 PowerPoint 表格转为 Markdown"""
    lines = []
    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(lines)


def _merge_content(text_blocks: list, images: list, tables: list) -> str:
    """合并所有内容"""
    parts = []
    
    for block in text_blocks:
        if block.get("content"):
            parts.append(block["content"])
    
    for img in images:
        if img.get("ocr_text"):
            parts.append(f"[图片内容] {img['ocr_text']}")
    
    for table in tables:
        if table.get("markdown"):
            parts.append(table["markdown"])
    
    return "\n\n".join(parts)
